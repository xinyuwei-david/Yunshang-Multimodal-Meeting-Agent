"""Generate renderer-neutral mind-map files and an editable PPTX."""

from __future__ import annotations

import json
import textwrap
from collections.abc import Callable
from html import escape
from pathlib import Path
from uuid import uuid4

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt

from .models import MeetingAnalysis, MindMapNode


def generate_artifacts(analysis: MeetingAnalysis, output_dir: Path) -> dict[str, Path]:
    output_dir.mkdir(parents=True, exist_ok=True)
    graph_json = output_dir / "mind-map.json"
    graph_svg = output_dir / "mind-map.svg"
    graph_png = output_dir / "mind-map.png"
    presentation = output_dir / "meeting-summary.pptx"
    analysis_json = output_dir / "meeting-analysis.json"

    _atomic_write_text(
        graph_json,
        json.dumps(analysis.mind_map.model_dump(), ensure_ascii=False, indent=2),
    )
    _atomic_write_text(analysis_json, analysis.model_dump_json(indent=2))
    _atomic_write_text(graph_svg, _mind_map_svg(analysis.mind_map))
    _atomic_generate(graph_png, lambda temporary: _mind_map_png(analysis.mind_map, temporary))
    _atomic_generate(presentation, lambda temporary: _presentation(analysis, temporary))
    return {
        "analysis": analysis_json,
        "mind_map_json": graph_json,
        "mind_map_svg": graph_svg,
        "mind_map_png": graph_png,
        "presentation": presentation,
    }


def _atomic_write_text(path: Path, text: str) -> None:
    temporary = _temporary_path(path)
    try:
        temporary.write_text(text, encoding="utf-8")
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _atomic_generate(path: Path, writer: Callable[[Path], None]) -> None:
    temporary = _temporary_path(path)
    try:
        writer(temporary)
        temporary.replace(path)
    finally:
        temporary.unlink(missing_ok=True)


def _temporary_path(path: Path) -> Path:
    return path.with_name(f".{path.name}.{uuid4().hex}.tmp")


def _mind_map_svg(root: MindMapNode) -> str:
    children = root.children[:6]
    height = max(480, 150 + len(children) * 110)
    lines = [
        '<svg xmlns="http://www.w3.org/2000/svg" width="1200" '
        f'height="{height}" viewBox="0 0 1200 {height}">',
        '<rect width="100%" height="100%" fill="#f7f9fc"/>',
        '<style>text{font-family:Segoe UI,Arial,sans-serif;fill:#1f2937}'
        '.root{font-size:25px;font-weight:700;fill:white}.node{font-size:18px;font-weight:600}'
        '.leaf{font-size:14px}</style>',
        '<rect x="430" y="28" width="340" height="72" rx="12" fill="#0f6cbd"/>',
        f'<text class="root" x="600" y="72" text-anchor="middle">{escape(root.label)}</text>',
    ]
    for index, child in enumerate(children):
        y = 145 + index * 110
        lines.extend(
            [
                f'<line x1="600" y1="100" x2="300" y2="{y + 28}" '
                'stroke="#0f6cbd" stroke-width="2"/>',
                f'<rect x="90" y="{y}" width="420" height="58" rx="10" '
                'fill="#e8f2fb" stroke="#0f6cbd"/>',
                f'<text class="node" x="110" y="{y + 35}">{escape(child.label)}</text>',
            ]
        )
        for leaf_index, leaf in enumerate(child.children[:3]):
            leaf_y = y + leaf_index * 26
            lines.append(
                f'<text class="leaf" x="680" y="{leaf_y + 20}">'
                f'- {escape(_clip(leaf.label, 68))}</text>'
            )
    lines.append("</svg>")
    return "\n".join(lines)


def _mind_map_png(root: MindMapNode, path: Path) -> None:
    image = Image.new("RGB", (1280, 720), "white")
    draw = ImageDraw.Draw(image)
    title_font = _font(22, bold=True)
    section_font = _font(18, bold=True)
    body_font = _font(15)
    draw.rounded_rectangle((440, 30, 840, 145), radius=18, fill="#0f6cbd")
    root_text = "\n".join(textwrap.wrap(root.label, width=34)[:3])
    draw.multiline_text(
        (640, 88),
        root_text,
        fill="white",
        font=title_font,
        anchor="mm",
        align="center",
        spacing=5,
    )
    for index, child in enumerate(root.children[:6]):
        left = 60 if index % 2 == 0 else 830
        top = 190 + (index // 2) * 170
        right = left + 390
        bottom = top + 130
        draw.line((640, 145, left + 195, top), fill="#0f6cbd", width=3)
        draw.rounded_rectangle(
            (left, top, right, bottom),
            radius=14,
            fill="#e8f2fb",
            outline="#0f6cbd",
            width=2,
        )
        draw.text(
            (left + 18, top + 15),
            _clip(child.label, 38),
            fill="#1f2937",
            font=section_font,
        )
        leaf_text = "\n".join(
            f"- {_clip(leaf.label, 42)}" for leaf in child.children[:3]
        )
        draw.multiline_text(
            (left + 18, top + 48),
            leaf_text,
            fill="#374151",
            font=body_font,
            spacing=6,
        )
    image.save(path, format="PNG")


def _font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    windows_name = "msyhbd.ttc" if bold else "msyh.ttc"
    noto_name = "NotoSansCJK-Bold.ttc" if bold else "NotoSansCJK-Regular.ttc"
    dejavu_name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    candidates = (
        Path("C:/Windows/Fonts") / windows_name,
        Path("/usr/share/fonts/opentype/noto") / noto_name,
        Path(dejavu_name),
    )
    for candidate in candidates:
        try:
            return ImageFont.truetype(str(candidate), size)
        except OSError:
            continue
    return ImageFont.load_default(size=size)


def _presentation(analysis: MeetingAnalysis, path: Path) -> None:
    presentation = Presentation()
    _title_slide(presentation, analysis.title, analysis.summary)
    _bullet_slide(presentation, "Topics", analysis.topics)
    _bullet_slide(presentation, "Decisions", analysis.decisions or ["No decision recorded"])
    action_lines = [
        " — ".join(filter(None, (item.description, item.owner, item.due)))
        for item in analysis.action_items
    ]
    _bullet_slide(presentation, "Action items", action_lines or ["No action item recorded"])
    _bullet_slide(
        presentation,
        "Open questions",
        analysis.open_questions or ["No open question recorded"],
    )
    presentation.save(path)


def _title_slide(presentation: Presentation, title: str, subtitle: str) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.25), Inches(8.5), Inches(1.2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(30)
    title_frame.paragraphs[0].font.bold = True
    subtitle_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(8.2), Inches(2.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(18)


def _bullet_slide(presentation: Presentation, title: str, items: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.55), Inches(8.7), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(26)
    title_frame.paragraphs[0].font.bold = True
    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(8.1), Inches(5.3))
    body_frame = body_box.text_frame
    body_frame.clear()
    for index, item in enumerate(items[:8]):
        paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
        paragraph.text = _clip(item, 180)
        paragraph.level = 0
        paragraph.font.size = Pt(19)
        paragraph.space_after = Pt(10)


def _clip(value: str, limit: int) -> str:
    compact = " ".join(textwrap.wrap(" ".join(value.split()), width=limit))
    return compact if len(compact) <= limit else compact[: limit - 1].rstrip() + "…"